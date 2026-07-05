"""Convert an uploaded document into Markdown.

Office formats (docx/pptx/xlsx) are converted into *structured* Markdown —
headings, lists and tables are preserved — so the header-aware chunker can
split along semantic boundaries. PDF has no reliable structure, so it falls back
to extracted plain text. Plain-text/markdown files are passed through as-is.
"""

from collections.abc import Iterator
from io import BytesIO
from pathlib import Path

from app.shared.exceptions.base import UnprocessableError

_TEXT_SUFFIXES = {".txt", ".md", ".markdown"}


class UnsupportedFormatError(UnprocessableError):
    """Raised when a document format cannot be parsed."""

    detail = "Unsupported document format"


# --------------------------------------------------------------------------- #
# PDF (plain text — PDFs carry no reliable heading/table structure)
# --------------------------------------------------------------------------- #
def _parse_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(data))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages).strip()


# --------------------------------------------------------------------------- #
# DOCX -> Markdown (headings, lists, tables), preserving document order
# --------------------------------------------------------------------------- #
def _docx_heading(style_name: str | None) -> str | None:
    if not style_name:
        return None
    if style_name == "Title":
        return "# "
    if style_name.startswith("Heading"):
        try:
            level = min(int(style_name.split()[-1]), 6)
        except ValueError:
            level = 1
        return "#" * level + " "
    return None


def _docx_list_marker(style_name: str | None) -> str | None:
    if not style_name:
        return None
    if "Number" in style_name:
        return "1. "
    if "Bullet" in style_name or style_name == "List Paragraph":
        return "- "
    return None


def _table_to_markdown(rows: list[list[str]]) -> str:
    rows = [r for r in rows if any(cell.strip() for cell in r)]
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    header = "| " + " | ".join(rows[0]) + " |"
    divider = "| " + " | ".join(["---"] * width) + " |"
    body = ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return "\n".join([header, divider, *body])


def _iter_docx_blocks(document: object) -> Iterator[str]:
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = document.element.body  # type: ignore[attr-defined]
    for child in body.iterchildren():
        tag = child.tag
        if tag.endswith("}p"):
            paragraph = Paragraph(child, document)
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = paragraph.style.name if paragraph.style else None
            heading = _docx_heading(style_name)
            if heading:
                yield f"{heading}{text}"
                continue
            marker = _docx_list_marker(style_name)
            yield f"{marker}{text}" if marker else text
        elif tag.endswith("}tbl"):
            table = Table(child, document)
            rows = [
                [cell.text.strip().replace("\n", " ") for cell in row.cells]
                for row in table.rows
            ]
            markdown = _table_to_markdown(rows)
            if markdown:
                yield markdown


def _parse_docx(data: bytes) -> str:
    from docx import Document

    document = Document(BytesIO(data))
    return "\n\n".join(_iter_docx_blocks(document)).strip()


# --------------------------------------------------------------------------- #
# PPTX -> Markdown (one section per slide)
# --------------------------------------------------------------------------- #
def _parse_pptx(data: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(BytesIO(data))
    blocks: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape else ""
        blocks.append(f"## Slide {index}: {title}".rstrip(": ").rstrip())
        for shape in slide.shapes:
            if shape is title_shape or not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    blocks.append(f"- {text}")
    return "\n\n".join(blocks).strip()


# --------------------------------------------------------------------------- #
# XLSX -> Markdown (one table per sheet)
# --------------------------------------------------------------------------- #
def _parse_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(BytesIO(data), read_only=True, data_only=True)
    blocks: list[str] = []
    for sheet in workbook.worksheets:
        rows = [
            ["" if cell is None else str(cell) for cell in row]
            for row in sheet.iter_rows(values_only=True)
        ]
        table = _table_to_markdown(rows)
        if table:
            blocks.append(f"## {sheet.title}\n\n{table}")
    workbook.close()
    return "\n\n".join(blocks).strip()


_PARSERS = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
    ".xlsx": _parse_xlsx,
}


def parse_to_markdown(filename: str, data: bytes) -> str:
    """Return the Markdown representation of a document."""
    suffix = Path(filename).suffix.lower()
    if suffix in _TEXT_SUFFIXES:
        return data.decode("utf-8", errors="replace").strip()
    parser = _PARSERS.get(suffix)
    if parser is not None:
        return parser(data)
    raise UnsupportedFormatError(f"Unsupported document format: {suffix}")
